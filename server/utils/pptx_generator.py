"""
PowerPoint presentation generator with template support.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from typing import List, Dict


def generate_presentation(title: str, slides: List[Dict], output_path: str):
    """
    Generate a PowerPoint presentation with various slide types.
    
    slides = [
        {"type": "title", "title": "Main Title", "subtitle": "Subtitle"},
        {"type": "content", "title": "Slide Title", "content": "Slide content text"},
        {"type": "bullets", "title": "Bullet Points", "bullets": ["Point 1", "Point 2"]},
        {"type": "chart", "title": "Chart", "chart_data": {...}},
    ]
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    for slide_data in slides:
        slide_type = slide_data.get("type", "content")
        
        if slide_type == "title":
            _add_title_slide(prs, slide_data)
        elif slide_type == "content":
            _add_content_slide(prs, slide_data)
        elif slide_type == "bullets":
            _add_bullet_slide(prs, slide_data)
        elif slide_type == "chart":
            _add_chart_slide(prs, slide_data)
        else:
            _add_content_slide(prs, slide_data)
    
    prs.save(output_path)


def _add_title_slide(prs: Presentation, data: Dict):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = data.get("title", "")
    subtitle.text = data.get("subtitle", "")
    
    # Style the title
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True


def _add_content_slide(prs: Presentation, data: Dict):
    """Add a content slide with title and body text."""
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = data.get("title", "")
    
    # Add content
    content = data.get("content", "")
    if len(slide.placeholders) > 1:
        body_shape = slide.placeholders[1]
        text_frame = body_shape.text_frame
        text_frame.text = content
    else:
        # If no placeholder, add a text box
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = txBox.text_frame
        text_frame.text = content


def _add_bullet_slide(prs: Presentation, data: Dict):
    """Add a slide with bullet points."""
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = data.get("title", "")
    
    # Add bullets
    bullets = data.get("bullets", [])
    if len(slide.placeholders) > 1:
        body_shape = slide.placeholders[1]
        text_frame = body_shape.text_frame
        text_frame.clear()
        
        for i, bullet_text in enumerate(bullets):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            p.text = bullet_text
            p.level = 0


def _add_chart_slide(prs: Presentation, data: Dict):
    """Add a slide with a chart."""
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = txBox.text_frame
    text_frame.text = data.get("title", "Chart")
    text_frame.paragraphs[0].font.size = Pt(32)
    text_frame.paragraphs[0].font.bold = True
    
    # Add chart
    chart_data_dict = data.get("chart_data", {})
    categories = chart_data_dict.get("categories", ["Category 1", "Category 2", "Category 3"])
    values = chart_data_dict.get("values", [10, 20, 30])
    
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series("Series 1", values)
    
    x, y, cx, cy = Inches(1.5), Inches(2), Inches(7), Inches(4.5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart


def generate_pptx(out_path, slides_text):
    """
    Legacy function for backward compatibility.
    Creates a simple presentation with text slides.
    """
    prs = Presentation()
    
    # Add a title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    title.text = "Generated Presentation"
    
    # Add content slides
    if isinstance(slides_text, str):
        slides_text = [slides_text]
    
    for i, text in enumerate(slides_text):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = f"Slide {i + 1}"
        body_shape = slide.placeholders[1]
        text_frame = body_shape.text_frame
        text_frame.text = text
    
    prs.save(out_path)
